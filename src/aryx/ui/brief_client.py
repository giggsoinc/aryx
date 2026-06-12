"""Brief-page client: extract text from an uploaded doc + call draft-brief.

Text extraction runs in the UI process (mirrors ingest_files) so the API
stays a clean JSON broker call. Supports PDF / DOCX / PPTX / TXT / MD.
"""
from __future__ import annotations

import io

from aryx.ui import api


def extract_text(uploaded) -> str:
    """Return plain text from a Streamlit UploadedFile (best-effort)."""
    if uploaded is None:
        return ""
    name = getattr(uploaded, "name", "").lower()
    data = uploaded.getvalue() if hasattr(uploaded, "getvalue") else uploaded.read()
    if name.endswith((".txt", ".md")):
        return data.decode("utf-8", errors="ignore")
    if name.endswith(".pdf"):
        return _pdf(data)
    if name.endswith((".docx", ".doc")):
        return _docx(data)
    if name.endswith((".pptx", ".ppt")):
        return _pptx(data)
    return ""


def _pdf(data: bytes) -> str:
    try:
        import fitz
        doc = fitz.open(stream=data, filetype="pdf")
        return "\n".join(page.get_text() for page in doc)
    except Exception:
        return ""


def _docx(data: bytes) -> str:
    try:
        from docx import Document
        doc = Document(io.BytesIO(data))
        return "\n".join(p.text for p in doc.paragraphs)
    except Exception:
        return ""


def _pptx(data: bytes) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
        lines = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    lines.append(shape.text_frame.text)
        return "\n".join(lines)
    except Exception:
        return ""


def draft_brief(seed: str, doc_text: str) -> dict:
    """Call the API to draft a 5-field brief from seed + document text."""
    wid = api.current_workspace()
    resp = api._post(f"/admin/workspaces/{wid}/draft-brief",
                     {"seed": seed, "doc_text": doc_text}, timeout=180)
    return resp.get("brief", {}) or {}

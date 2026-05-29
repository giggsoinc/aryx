"""PPTX document connector: slide text + embedded image OCR (Inc 8)."""
from __future__ import annotations

import io
import logging
from collections.abc import Iterator
from pathlib import Path

logger = logging.getLogger(__name__)


class PptxConnector:
    """Extracts text from PPTX slides via python-pptx; OCRs embedded images."""

    def __init__(self, path: Path) -> None:
        self._path = path

    def extract_pages(self) -> Iterator[tuple[int, str]]:
        """Yield (slide_num, text) for every slide, including OCR of embedded images."""
        from pptx import Presentation

        prs = Presentation(str(self._path))
        for slide_num, slide in enumerate(prs.slides, start=1):
            parts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    parts.append(shape.text_frame.text)
                if hasattr(shape, "image"):
                    ocr_text = _ocr_image_bytes(shape.image.blob)
                    if ocr_text:
                        parts.append(ocr_text)
            yield slide_num, "\n".join(parts)

    def close(self) -> None:
        pass


def _ocr_image_bytes(img_bytes: bytes) -> str:
    try:
        import pytesseract
        from PIL import Image
        img = Image.open(io.BytesIO(img_bytes))
        return pytesseract.image_to_string(img).strip()
    except Exception as exc:
        logger.debug("ocr skipped: %s", exc)
        return ""
